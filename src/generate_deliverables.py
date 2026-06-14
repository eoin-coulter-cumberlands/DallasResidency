"""
Generate the final submission artifacts for the e-commerce analytics project.

This script turns the cleaned datasets and analysis helpers into a reproducible
submission package:

- dashboard HTML
- executive summary markdown
- technical report DOCX
- presentation PPTX
- chart PNGs
- supporting CSV tables

Run:
    python3 src/generate_deliverables.py
"""

from __future__ import annotations

import base64
import json
import os
import sys
from pathlib import Path

ROOT = Path(__file__).resolve().parents[1]
if str(ROOT) not in sys.path:
    sys.path.insert(0, str(ROOT))

os.environ.setdefault("MPLCONFIGDIR", str(ROOT / ".matplotlib"))

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import numpy as np
import pandas as pd
import plotly.graph_objects as go
import seaborn as sns
from docx import Document
from docx.enum.section import WD_SECTION
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from docx.shared import Inches, Pt, RGBColor as DocxRGBColor
from pptx import Presentation
from pptx.dml.color import RGBColor as PptRGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches as PptInches
from pptx.util import Pt as PptPt
from plotly.offline import get_plotlyjs
from plotly.utils import PlotlyJSONEncoder

from src import analysis_helpers as ah
DELIVERABLES = ROOT / "deliverables"
FIGURES = DELIVERABLES / "figures"
TABLES = DELIVERABLES / "tables"
DASHBOARD = DELIVERABLES / "dashboard"
REPORTS = DELIVERABLES / "report"
PRESENTATION = DELIVERABLES / "presentation"

BRAND = {
    "ink": "#16233B",
    "copper": "#C96A3D",
    "saffron": "#E3A72F",
    "sage": "#6E8B74",
    "rose": "#B5525C",
    "mist": "#F4EFE8",
    "stone": "#E4DDD3",
    "text": "#1C2430",
}


def ensure_dirs() -> None:
    for path in [DELIVERABLES, FIGURES, TABLES, DASHBOARD, REPORTS, PRESENTATION]:
        path.mkdir(parents=True, exist_ok=True)


def setup_style() -> None:
    sns.set_theme(style="whitegrid")
    plt.rcParams.update(
        {
            "figure.figsize": (12, 7),
            "axes.titlesize": 18,
            "axes.labelsize": 12,
            "xtick.labelsize": 10,
            "ytick.labelsize": 10,
        }
    )


def save_plot(fig: plt.Figure, name: str) -> Path:
    path = FIGURES / name
    fig.tight_layout()
    fig.savefig(path, dpi=180, bbox_inches="tight")
    plt.close(fig)
    return path


def save_table(df: pd.DataFrame, name: str) -> Path:
    path = TABLES / name
    df.to_csv(path, index=False)
    return path


def plotly_json(fig: go.Figure) -> str:
    return json.dumps(fig, cls=PlotlyJSONEncoder)


def format_currency(value: float) -> str:
    return f"GBP {value:,.0f}"


def format_pct(value: float) -> str:
    return f"{value * 100:.1f}%"


def card_html(title: str, value: str, subtitle: str) -> str:
    return f"""
    <div class="card">
      <div class="card-title">{title}</div>
      <div class="card-value">{value}</div>
      <div class="card-subtitle">{subtitle}</div>
    </div>
    """


def image_to_data_uri(path: Path) -> str:
    encoded = base64.b64encode(path.read_bytes()).decode("ascii")
    return f"data:image/png;base64,{encoded}"


def add_docx_page_number(paragraph) -> None:
    run = paragraph.add_run()
    fld_char_begin = OxmlElement("w:fldChar")
    fld_char_begin.set(qn("w:fldCharType"), "begin")
    instr_text = OxmlElement("w:instrText")
    instr_text.set(qn("xml:space"), "preserve")
    instr_text.text = "PAGE"
    fld_char_end = OxmlElement("w:fldChar")
    fld_char_end.set(qn("w:fldCharType"), "end")
    run._r.append(fld_char_begin)
    run._r.append(instr_text)
    run._r.append(fld_char_end)


def set_doc_defaults(doc: Document) -> None:
    section = doc.sections[0]
    section.top_margin = Inches(0.8)
    section.bottom_margin = Inches(0.7)
    section.left_margin = Inches(0.8)
    section.right_margin = Inches(0.8)

    styles = doc.styles
    styles["Normal"].font.name = "Arial"
    styles["Normal"].font.size = Pt(10.5)

    for style_name, size, bold in [
        ("Title", 24, True),
        ("Heading 1", 16, True),
        ("Heading 2", 13, True),
    ]:
        style = styles[style_name]
        style.font.name = "Arial"
        style.font.size = Pt(size)
        style.font.bold = bold

    footer = section.footer.paragraphs[0]
    footer.alignment = WD_ALIGN_PARAGRAPH.RIGHT
    add_docx_page_number(footer)


def add_doc_title(doc: Document, title: str, subtitle: str) -> None:
    p = doc.add_paragraph()
    p.style = doc.styles["Title"]
    p.alignment = WD_ALIGN_PARAGRAPH.LEFT
    run = p.add_run(title)
    run.font.color.rgb = DocxRGBColor(20, 58, 82)
    sub = doc.add_paragraph()
    sub.alignment = WD_ALIGN_PARAGRAPH.LEFT
    sub_run = sub.add_run(subtitle)
    sub_run.font.name = "Arial"
    sub_run.font.size = Pt(11)
    sub_run.italic = True


def add_bullets(doc: Document, bullets: list[str]) -> None:
    for bullet in bullets:
        p = doc.add_paragraph(style="List Bullet")
        p.add_run(bullet)


def add_picture(doc: Document, path: Path, width: float = 6.5) -> None:
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p.add_run().add_picture(str(path), width=Inches(width))


def build_analysis_payload() -> dict:
    headline = ah.headline_metrics()
    monthly = ah.monthly_revenue()
    products_rev = ah.top_products_by_revenue(10)
    products_qty = ah.top_products_by_quantity(10)
    customers = ah.top_customers_by_clv(10)
    rfm = ah.rfm_segments()
    countries = ah.revenue_by_country()
    day = ah.revenue_by_day_of_week()
    hour = ah.revenue_by_hour()
    basket = ah.basket_summary()
    basket_metrics = ah.basket_metrics()
    price_bands = ah.revenue_by_price_band()
    returns = ah.returns_summary()
    top_returns = ah.top_returned_products(10)
    returns_time = ah.returns_over_time()
    cohorts = ah.cohort_retention_rate()
    forecast = ah.revenue_forecast(3)
    repeat = ah.repeat_vs_onetime()

    top_month = monthly.loc[monthly["Revenue"].idxmax()]
    best_day = day.dropna().loc[day["Revenue"].idxmax()]
    best_hour = hour.loc[hour["Revenue"].idxmax()]

    uk_share = countries.loc[countries["Country"] == "United Kingdom", "RevenueShare"].iat[0]
    international_share = 1 - uk_share
    repeat_customers = int(repeat.loc[repeat["Segment"] == "Repeat", "Customers"].iat[0])
    one_time_customers = int(repeat.loc[repeat["Segment"] == "One-time", "Customers"].iat[0])
    repeat_share = repeat_customers / (repeat_customers + one_time_customers)
    return_value_ratio = abs(returns["financial_impact"]) / headline["total_revenue"]

    findings = [
        f"Revenue reached {format_currency(headline['total_revenue'])} from {headline['total_orders']:,} orders, with an average order value of {format_currency(headline['avg_order_value'])}.",
        f"Peak trading month was {top_month['YearMonth']} at {format_currency(top_month['Revenue'])}, and September to November formed the strongest seasonal run-up into year-end.",
        f"The business is highly concentrated in the UK ({format_pct(uk_share)} of revenue), but the Netherlands, EIRE, Germany, and France form a credible international expansion shortlist.",
        f"Repeat customers account for {format_pct(repeat_share)} of identified customers, while the December 2010 cohort retained {cohorts.loc['2010-12', 11] * 100:.1f}% of customers by month 11.",
        f"Returns and cancellations reduced value by {format_currency(abs(returns['financial_impact']))}, or about {format_pct(return_value_ratio)} of clean revenue.",
        f"The highest-value promotional window is late morning through mid-afternoon, especially Thursdays at {int(best_hour['Hour']):02d}:00 and Tuesdays/Thursdays overall.",
    ]

    recommendations = [
        "Defend and scale top giftable SKUs such as REGENCY CAKESTAND 3 TIER, PARTY BUNTING, and JUMBO BAG RED RETROSPOT with priority inventory and featured placement.",
        "Target high-value repeat buyers with replenishment and seasonal launch campaigns 30-45 days after purchase to lift cohort retention.",
        "Concentrate promotional sends on Tuesday through Thursday and schedule merchandising pushes between 10:00 and 15:00 when conversion value is strongest.",
        "Audit the unusually large return lines and manual adjustment codes separately so operational noise does not hide true merchandise returns.",
        "Prioritize expansion tests in the Netherlands, EIRE, Germany, and France before broader international rollout because they already contribute meaningful revenue at low customer counts.",
    ]

    return {
        "headline": headline,
        "monthly": monthly,
        "products_rev": products_rev,
        "products_qty": products_qty,
        "customers": customers,
        "rfm": rfm,
        "countries": countries,
        "day": day,
        "hour": hour,
        "basket": basket,
        "basket_metrics": basket_metrics,
        "price_bands": price_bands,
        "returns": returns,
        "top_returns": top_returns,
        "returns_time": returns_time,
        "cohorts": cohorts,
        "forecast": forecast,
        "repeat": repeat,
        "findings": findings,
        "recommendations": recommendations,
        "summary_bits": {
            "top_month": top_month["YearMonth"],
            "top_month_revenue": float(top_month["Revenue"]),
            "best_day": best_day["DayOfWeek"],
            "best_hour": int(best_hour["Hour"]),
            "uk_share": uk_share,
            "international_share": international_share,
            "repeat_share": repeat_share,
            "return_value_ratio": return_value_ratio,
        },
    }


def build_figures(payload: dict) -> dict[str, Path]:
    monthly = payload["monthly"].copy()
    products_rev = payload["products_rev"].copy().sort_values("TotalRevenue")
    rfm_counts = payload["rfm"]["Segment"].value_counts().reset_index()
    rfm_counts.columns = ["Segment", "Customers"]
    countries = payload["countries"].head(10).sort_values("Revenue")
    day = payload["day"].dropna()
    hour = payload["hour"]
    basket = payload["basket"]
    price_bands = payload["price_bands"]
    top_returns = payload["top_returns"].copy().sort_values("UnitsReturned")
    cohorts = payload["cohorts"].fillna(0)
    forecast = payload["forecast"]

    files = {}

    fig, ax = plt.subplots()
    ax.plot(monthly["YearMonth"], monthly["Revenue"], marker="o", linewidth=2.5, color=BRAND["ink"])
    ax.set_title("Monthly Revenue Trend")
    ax.set_xlabel("Month")
    ax.set_ylabel("Revenue (GBP)")
    ax.tick_params(axis="x", rotation=45)
    files["monthly_revenue"] = save_plot(fig, "monthly_revenue.png")

    fig, ax = plt.subplots()
    ax.barh(products_rev["Description"], products_rev["TotalRevenue"], color=BRAND["copper"])
    ax.set_title("Top 10 Products by Revenue")
    ax.set_xlabel("Revenue (GBP)")
    files["top_products_revenue"] = save_plot(fig, "top_products_revenue.png")

    fig, ax = plt.subplots()
    ax.bar(rfm_counts["Segment"], rfm_counts["Customers"], color=[BRAND["rose"], BRAND["ink"], BRAND["sage"], BRAND["saffron"], "#9A8F7A"])
    ax.set_title("Customer Segments (RFM)")
    ax.set_ylabel("Customers")
    ax.tick_params(axis="x", rotation=20)
    files["rfm_segments"] = save_plot(fig, "rfm_segments.png")

    fig, ax = plt.subplots()
    ax.barh(countries["Country"], countries["Revenue"], color=BRAND["sage"])
    ax.set_title("Top 10 Countries by Revenue")
    ax.set_xlabel("Revenue (GBP)")
    files["top_countries"] = save_plot(fig, "top_countries.png")

    fig, ax = plt.subplots()
    ax.bar(day["DayOfWeek"], day["Revenue"], color=BRAND["saffron"])
    ax.set_title("Revenue by Day of Week")
    ax.set_ylabel("Revenue (GBP)")
    ax.tick_params(axis="x", rotation=20)
    files["day_of_week"] = save_plot(fig, "day_of_week.png")

    fig, ax = plt.subplots()
    ax.plot(hour["Hour"], hour["Revenue"], marker="o", color=BRAND["copper"], linewidth=2.5)
    ax.set_title("Revenue by Hour")
    ax.set_xlabel("Hour")
    ax.set_ylabel("Revenue (GBP)")
    files["hour_of_day"] = save_plot(fig, "hour_of_day.png")

    fig, ax = plt.subplots()
    sns.histplot(basket["Value"].clip(upper=basket["Value"].quantile(0.95)), bins=40, color=BRAND["ink"], ax=ax)
    ax.set_title("Basket Value Distribution (trimmed at 95th percentile)")
    ax.set_xlabel("Order Value (GBP)")
    files["basket_distribution"] = save_plot(fig, "basket_distribution.png")

    fig, ax = plt.subplots()
    ax.bar(price_bands["PriceBand"].astype(str), price_bands["Revenue"], color=BRAND["sage"])
    ax.set_title("Revenue by Price Band")
    ax.set_ylabel("Revenue (GBP)")
    ax.tick_params(axis="x", rotation=25)
    files["price_bands"] = save_plot(fig, "price_bands.png")

    fig, ax = plt.subplots()
    ax.barh(top_returns["Description"], top_returns["UnitsReturned"], color=BRAND["rose"])
    ax.set_title("Most Returned Products")
    ax.set_xlabel("Units Returned")
    files["top_returns"] = save_plot(fig, "top_returns.png")

    fig, ax = plt.subplots(figsize=(12, 6))
    sns.heatmap(cohorts, cmap="YlGnBu", annot=True, fmt=".0%", cbar_kws={"format": "%.0f%%"}, ax=ax)
    ax.set_title("Monthly Cohort Retention")
    ax.set_xlabel("Months Since First Purchase")
    ax.set_ylabel("Cohort Month")
    files["cohort_heatmap"] = save_plot(fig, "cohort_heatmap.png")

    fig, ax = plt.subplots()
    actual = forecast[forecast["Type"] == "Actual"]
    pred = forecast[forecast["Type"] == "Forecast"]
    ax.plot(actual["YearMonth"], actual["Revenue"], marker="o", color=BRAND["ink"], linewidth=2.5, label="Actual")
    ax.plot(pred["YearMonth"], pred["Revenue"], marker="o", linestyle="--", color=BRAND["copper"], linewidth=2.5, label="Forecast")
    ax.set_title("Revenue Forecast")
    ax.set_ylabel("Revenue (GBP)")
    ax.tick_params(axis="x", rotation=45)
    ax.legend()
    files["forecast"] = save_plot(fig, "forecast.png")

    return files


def build_tables(payload: dict) -> None:
    save_table(payload["monthly"], "monthly_revenue.csv")
    save_table(payload["products_rev"], "top_products_revenue.csv")
    save_table(payload["products_qty"], "top_products_quantity.csv")
    save_table(payload["customers"], "top_customers.csv")
    save_table(payload["countries"], "revenue_by_country.csv")
    save_table(payload["price_bands"], "revenue_by_price_band.csv")
    save_table(payload["top_returns"], "top_returned_products.csv")
    save_table(payload["cohorts"].reset_index(), "cohort_retention.csv")
    save_table(payload["forecast"], "revenue_forecast.csv")


def build_dashboard(payload: dict, figures: dict[str, Path]) -> Path:
    headline = payload["headline"]
    bits = payload["summary_bits"]
    cards = "\n".join(
        [
            card_html("Revenue", format_currency(headline["total_revenue"]), f"{headline['total_orders']:,} orders"),
            card_html("Avg Order Value", format_currency(headline["avg_order_value"]), f"Median {format_currency(headline['median_order_value'])}"),
            card_html("Customers", f"{headline['total_customers']:,}", f"Across {headline['countries']} countries"),
            card_html("Returns Impact", format_pct(bits["return_value_ratio"]), format_currency(abs(payload["returns"]["financial_impact"]))),
        ]
    )

    finding_items = "\n".join(f"<li>{item}</li>" for item in payload["findings"])
    recommendation_items = "\n".join(f"<li>{item}</li>" for item in payload["recommendations"])

    top_products_html = (
        payload["products_rev"][["Description", "TotalRevenue", "TotalQuantity"]]
        .assign(TotalRevenue=lambda d: d["TotalRevenue"].map(lambda x: f"{x:,.0f}"))
        .to_html(index=False, classes="table", border=0)
    )
    countries_html = (
        payload["countries"].head(8)[["Country", "Revenue", "Orders"]]
        .assign(Revenue=lambda d: d["Revenue"].map(lambda x: f"{x:,.0f}"))
        .to_html(index=False, classes="table", border=0)
    )
    rfm_counts = payload["rfm"]["Segment"].value_counts()
    countries_top = payload["countries"].head(10).sort_values("Revenue", ascending=True)
    top_products = payload["products_rev"].sort_values("TotalRevenue", ascending=True)
    top_returns = payload["top_returns"].sort_values("UnitsReturned", ascending=True)
    price_bands = payload["price_bands"].copy()
    monthly = payload["monthly"].copy()
    hour = payload["hour"].copy()
    basket = payload["basket"].copy()
    cohorts = payload["cohorts"].fillna(0)
    forecast = payload["forecast"].copy()

    material_layout = dict(
        paper_bgcolor="rgba(0,0,0,0)",
        plot_bgcolor="#FFFFFF",
        margin=dict(l=40, r=24, t=48, b=42),
        font=dict(family="Arial, sans-serif", color="#1F2933", size=12),
        title=dict(font=dict(size=18, color=BRAND["ink"])),
        hoverlabel=dict(bgcolor=BRAND["ink"], font=dict(color="#FFFFFF")),
    )

    chart_specs = [
        {
            "id": "monthlyRevenue",
            "title": "Monthly Revenue Trend",
            "insight": "Strong Q4 acceleration with a November peak and a partial December close.",
            "figure": go.Figure(
                data=[
                    go.Scatter(
                        x=monthly["YearMonth"],
                        y=monthly["Revenue"],
                        mode="lines+markers",
                        line=dict(color=BRAND["ink"], width=3),
                        marker=dict(size=9, color=BRAND["saffron"]),
                        hovertemplate="Month %{x}<br>Revenue GBP %{y:,.0f}<extra></extra>",
                    )
                ],
                layout=go.Layout(
                    **material_layout,
                    xaxis=dict(title="Month", tickangle=-35, showgrid=False),
                    yaxis=dict(title="Revenue (GBP)", gridcolor="#E6ECF2"),
                ),
            ),
        },
        {
            "id": "topProducts",
            "title": "Top Products by Revenue",
            "insight": "Hero SKUs create a disproportionate share of value.",
            "figure": go.Figure(
                data=[
                    go.Bar(
                        x=top_products["TotalRevenue"],
                        y=top_products["Description"],
                        orientation="h",
                        marker=dict(color=BRAND["copper"]),
                        customdata=top_products["TotalQuantity"],
                        hovertemplate="%{y}<br>Revenue GBP %{x:,.0f}<br>Units %{customdata:,.0f}<extra></extra>",
                    )
                ],
                layout=go.Layout(
                    **material_layout,
                    xaxis=dict(title="Revenue (GBP)", gridcolor="#E6ECF2"),
                    yaxis=dict(title=""),
                ),
            ),
        },
        {
            "id": "customerSegments",
            "title": "Customer Segments",
            "insight": "Champions are large, but lapse recovery is the biggest unlocked lever.",
            "figure": go.Figure(
                data=[
                    go.Bar(
                        x=rfm_counts.index.tolist(),
                        y=rfm_counts.values.tolist(),
                        marker=dict(color=[BRAND["rose"], BRAND["ink"], BRAND["sage"], BRAND["saffron"], "#9A8F7A"]),
                        hovertemplate="%{x}<br>Customers %{y:,.0f}<extra></extra>",
                    )
                ],
                layout=go.Layout(
                    **material_layout,
                    xaxis=dict(title="", tickangle=-15, showgrid=False),
                    yaxis=dict(title="Customers", gridcolor="#E6ECF2"),
                ),
            ),
        },
        {
            "id": "countryRevenue",
            "title": "Top Countries",
            "insight": "International demand exists already, but it is concentrated in a few markets.",
            "figure": go.Figure(
                data=[
                    go.Bar(
                        x=countries_top["Revenue"],
                        y=countries_top["Country"],
                        orientation="h",
                        marker=dict(color=BRAND["sage"]),
                        customdata=np.stack([countries_top["Orders"], countries_top["Customers"]], axis=-1),
                        hovertemplate="%{y}<br>Revenue GBP %{x:,.0f}<br>Orders %{customdata[0]:,.0f}<br>Customers %{customdata[1]:,.0f}<extra></extra>",
                    )
                ],
                layout=go.Layout(
                    **material_layout,
                    xaxis=dict(title="Revenue (GBP)", gridcolor="#E6ECF2"),
                    yaxis=dict(title=""),
                ),
            ),
        },
        {
            "id": "dayPattern",
            "title": "Revenue by Day of Week",
            "insight": "Midweek outperforms weekends by a wide margin.",
            "figure": go.Figure(
                data=[
                    go.Bar(
                        x=payload["day"]["DayOfWeek"].tolist(),
                        y=payload["day"]["Revenue"].fillna(0).tolist(),
                        marker=dict(color=BRAND["saffron"]),
                        hovertemplate="%{x}<br>Revenue GBP %{y:,.0f}<extra></extra>",
                    )
                ],
                layout=go.Layout(
                    **material_layout,
                    xaxis=dict(title="", tickangle=-20, showgrid=False),
                    yaxis=dict(title="Revenue (GBP)", gridcolor="#E6ECF2"),
                ),
            ),
        },
        {
            "id": "hourPattern",
            "title": "Revenue by Hour",
            "insight": "Late morning through mid-afternoon is the highest-value window.",
            "figure": go.Figure(
                data=[
                    go.Scatter(
                        x=hour["Hour"],
                        y=hour["Revenue"],
                        mode="lines+markers",
                        line=dict(color=BRAND["copper"], width=3),
                        marker=dict(color=BRAND["ink"], size=8),
                        hovertemplate="%{x}:00<br>Revenue GBP %{y:,.0f}<extra></extra>",
                    )
                ],
                layout=go.Layout(
                    **material_layout,
                    xaxis=dict(title="Hour", dtick=1, showgrid=False),
                    yaxis=dict(title="Revenue (GBP)", gridcolor="#E6ECF2"),
                ),
            ),
        },
        {
            "id": "basketDistribution",
            "title": "Basket Value Distribution",
            "insight": "Most orders are moderate, with a long premium tail pulling the average upward.",
            "figure": go.Figure(
                data=[
                    go.Histogram(
                        x=basket["Value"].clip(upper=basket["Value"].quantile(0.95)),
                        marker=dict(color=BRAND["ink"]),
                        nbinsx=35,
                        hovertemplate="Order value GBP %{x:,.0f}<br>Count %{y}<extra></extra>",
                    )
                ],
                layout=go.Layout(
                    **material_layout,
                    xaxis=dict(title="Order Value (GBP)", gridcolor="#E6ECF2"),
                    yaxis=dict(title="Orders", gridcolor="#E6ECF2"),
                    bargap=0.05,
                ),
            ),
        },
        {
            "id": "priceBands",
            "title": "Revenue by Price Band",
            "insight": "The GBP 2-5 zone is the commercial center of gravity.",
            "figure": go.Figure(
                data=[
                    go.Bar(
                        x=price_bands["PriceBand"].astype(str),
                        y=price_bands["Revenue"],
                        marker=dict(color=BRAND["sage"]),
                        customdata=np.stack([price_bands["Units"], price_bands["Lines"]], axis=-1),
                        hovertemplate="%{x}<br>Revenue GBP %{y:,.0f}<br>Units %{customdata[0]:,.0f}<br>Lines %{customdata[1]:,.0f}<extra></extra>",
                    )
                ],
                layout=go.Layout(
                    **material_layout,
                    xaxis=dict(title="", tickangle=-20, showgrid=False),
                    yaxis=dict(title="Revenue (GBP)", gridcolor="#E6ECF2"),
                ),
            ),
        },
        {
            "id": "returns",
            "title": "Most Returned Products",
            "insight": "A few outlier codes dominate the return story and need root-cause review.",
            "figure": go.Figure(
                data=[
                    go.Bar(
                        x=top_returns["UnitsReturned"],
                        y=top_returns["Description"],
                        orientation="h",
                        marker=dict(color=BRAND["rose"]),
                        customdata=top_returns["ValueReturned"],
                        hovertemplate="%{y}<br>Units Returned %{x:,.0f}<br>Value GBP %{customdata:,.0f}<extra></extra>",
                    )
                ],
                layout=go.Layout(
                    **material_layout,
                    xaxis=dict(title="Units Returned", gridcolor="#E6ECF2"),
                    yaxis=dict(title=""),
                ),
            ),
        },
        {
            "id": "cohorts",
            "title": "Cohort Retention",
            "insight": "Retention softens quickly after acquisition, then stabilizes for stronger cohorts.",
            "figure": go.Figure(
                data=[
                    go.Heatmap(
                        z=cohorts.values,
                        x=[str(col) for col in cohorts.columns],
                        y=cohorts.index.tolist(),
                        colorscale="Tealgrn",
                        colorbar=dict(title="Rate"),
                        hovertemplate="Cohort %{y}<br>Month %{x}<br>Retention %{z:.1%}<extra></extra>",
                    )
                ],
                layout=go.Layout(
                    **material_layout,
                    xaxis=dict(title="Months Since First Purchase"),
                    yaxis=dict(title="Cohort Month"),
                ),
            ),
        },
        {
            "id": "forecast",
            "title": "Revenue Forecast",
            "insight": "The baseline trend stays elevated into early 2012, even before richer modeling.",
            "figure": go.Figure(
                data=[
                    go.Scatter(
                        x=forecast.loc[forecast["Type"] == "Actual", "YearMonth"],
                        y=forecast.loc[forecast["Type"] == "Actual", "Revenue"],
                        mode="lines+markers",
                        name="Actual",
                        line=dict(color=BRAND["ink"], width=3),
                        marker=dict(size=8, color=BRAND["saffron"]),
                    ),
                    go.Scatter(
                        x=forecast.loc[forecast["Type"] == "Forecast", "YearMonth"],
                        y=forecast.loc[forecast["Type"] == "Forecast", "Revenue"],
                        mode="lines+markers",
                        name="Forecast",
                        line=dict(color=BRAND["copper"], width=3, dash="dash"),
                        marker=dict(size=8, color=BRAND["copper"]),
                    ),
                ],
                layout=go.Layout(
                    **material_layout,
                    xaxis=dict(title="Month", tickangle=-35, showgrid=False),
                    yaxis=dict(title="Revenue (GBP)", gridcolor="#E6ECF2"),
                    legend=dict(orientation="h", yanchor="bottom", y=1.02, xanchor="right", x=1),
                ),
            ),
        },
    ]

    chart_blocks = "\n".join(
        f"""
        <section class="chart-card interactive-card" data-chart-id="{spec['id']}">
          <div class="chart-card-top">
            <div>
              <h3>{spec['title']}</h3>
              <p>{spec['insight']}</p>
            </div>
            <button class="expand-button" type="button" data-chart-id="{spec['id']}">Expand</button>
          </div>
          <div class="chart-shell">
            <div class="chart-viewport" id="{spec['id']}"></div>
          </div>
        </section>
        """
        for spec in chart_specs
    )

    chart_json = "{\n" + ",\n".join(
        f'  "{spec["id"]}": {plotly_json(spec["figure"])}'
        for spec in chart_specs
    ) + "\n}"
    plotly_js = get_plotlyjs()

    html = f"""<!DOCTYPE html>
<html lang="en">
<head>
  <meta charset="utf-8">
  <meta name="viewport" content="width=device-width, initial-scale=1">
  <title>E-Commerce Business Analytics Dashboard</title>
  <style>
    :root {{
      --md-primary: #425d4a;
      --md-primary-dark: #16233b;
      --md-primary-soft: #e8dfd2;
      --md-secondary: #c96a3d;
      --md-accent: #e3a72f;
      --md-surface: #fffdf9;
      --md-surface-2: #f7f1ea;
      --md-outline: #ddd3c7;
      --md-text: #1c2430;
      --md-muted: #6c655d;
      --md-danger: #b5525c;
      --md-shadow: 0 16px 34px rgba(22, 35, 59, 0.10), 0 3px 10px rgba(22, 35, 59, 0.05);
    }}
    * {{ box-sizing: border-box; }}
    body {{
      margin: 0;
      font-family: Arial, sans-serif;
      color: var(--md-text);
      background:
        radial-gradient(circle at top right, rgba(201, 106, 61, 0.20), transparent 24%),
        radial-gradient(circle at top left, rgba(110, 139, 116, 0.18), transparent 20%),
        linear-gradient(180deg, #fbf7f1 0%, #efe7dc 100%);
    }}
    .app-bar {{
      position: sticky;
      top: 0;
      z-index: 15;
      display: flex;
      align-items: center;
      justify-content: space-between;
      padding: 16px 6vw;
      background: rgba(255,255,255,0.84);
      backdrop-filter: blur(16px);
      border-bottom: 1px solid rgba(221, 211, 199, 0.9);
    }}
    .brand {{
      font-size: 1rem;
      font-weight: 700;
      letter-spacing: 0.04em;
      color: var(--md-primary-dark);
      text-transform: uppercase;
    }}
    .chip-row {{
      display: flex;
      gap: 10px;
      flex-wrap: wrap;
    }}
    .chip {{
      padding: 8px 12px;
      border-radius: 999px;
      background: var(--md-primary-soft);
      color: var(--md-primary-dark);
      font-size: 0.86rem;
      font-weight: 600;
    }}
    header {{
      padding: 34px 6vw 28px;
    }}
    .hero {{
      display: grid;
      grid-template-columns: 1.2fr 0.8fr;
      gap: 24px;
      align-items: stretch;
    }}
    .hero-panel, .surface {{
      background: rgba(255,255,255,0.86);
      border: 1px solid rgba(215, 225, 231, 0.85);
      border-radius: 28px;
      box-shadow: var(--md-shadow);
    }}
    .hero-copy {{
      padding: 34px 34px 28px;
      background:
        linear-gradient(135deg, var(--md-primary-dark), var(--md-primary)),
        linear-gradient(180deg, rgba(255,255,255,0.08), transparent);
      color: #fff;
    }}
    .hero-copy h1 {{
      margin: 0 0 12px;
      font-size: 2.8rem;
      letter-spacing: 0.02em;
    }}
    .hero-copy p {{
      line-height: 1.6;
      font-size: 1.05rem;
      max-width: 760px;
      color: rgba(255,255,255,0.86);
    }}
    .hero-side {{
      padding: 26px;
      display: flex;
      flex-direction: column;
      justify-content: space-between;
      background: linear-gradient(180deg, #fffdf9, #f6efe6);
    }}
    .hero-side h2 {{
      margin: 0 0 10px;
      color: var(--md-primary-dark);
      font-size: 1.1rem;
    }}
    .hero-side p {{
      margin: 0;
      color: var(--md-muted);
      line-height: 1.6;
    }}
    main {{
      padding: 10px 6vw 56px;
    }}
    .cards {{
      display: grid;
      grid-template-columns: repeat(auto-fit, minmax(220px, 1fr));
      gap: 18px;
      margin: 26px 0 28px;
    }}
    .card, .panel, .chart-card {{
      background: var(--md-surface);
      border-radius: 24px;
      border: 1px solid var(--md-outline);
      box-shadow: var(--md-shadow);
      padding: 22px;
    }}
    .card-title {{
      font-size: 0.85rem;
      text-transform: uppercase;
      color: var(--md-muted);
      letter-spacing: 0.08em;
    }}
    .card-value {{
      margin: 8px 0 6px;
      font-size: 1.8rem;
      font-weight: 700;
      color: var(--md-primary-dark);
    }}
    .card-subtitle {{
      color: var(--md-muted);
      font-size: 0.95rem;
    }}
    .two-col {{
      display: grid;
      grid-template-columns: 1.1fr 0.9fr;
      gap: 20px;
      margin-bottom: 22px;
    }}
    .panel h2, .panel h3, .chart-card h3 {{
      margin-top: 0;
      color: var(--md-primary-dark);
    }}
    .panel p {{
      color: var(--md-muted);
      line-height: 1.6;
    }}
    .table {{
      width: 100%;
      border-collapse: collapse;
      margin-top: 10px;
      font-size: 0.95rem;
    }}
    .table th, .table td {{
      padding: 10px 12px;
      border-bottom: 1px solid var(--md-outline);
      text-align: left;
    }}
    .table th {{
      background: var(--md-surface-2);
    }}
    .chart-grid {{
      display: grid;
      grid-template-columns: repeat(auto-fit, minmax(320px, 1fr));
      gap: 20px;
    }}
    .chart-card {{
      padding-bottom: 18px;
      transition: transform 0.2s ease, box-shadow 0.2s ease;
      cursor: pointer;
    }}
    .chart-card:hover {{
      transform: translateY(-2px);
      box-shadow: 0 18px 36px rgba(22, 35, 59, 0.14), 0 4px 10px rgba(22, 35, 59, 0.08);
    }}
    .chart-card-top {{
      display: flex;
      align-items: flex-start;
      justify-content: space-between;
      gap: 16px;
      margin-bottom: 12px;
    }}
    .chart-card-top p {{
      margin: 4px 0 0;
      color: var(--md-muted);
      font-size: 0.95rem;
      line-height: 1.5;
    }}
    .expand-button {{
      flex: 0 0 auto;
      border: none;
      border-radius: 999px;
      background: var(--md-primary-soft);
      color: var(--md-primary-dark);
      font-weight: 700;
      padding: 10px 14px;
      cursor: pointer;
    }}
    .chart-shell {{
      border-radius: 18px;
      background: linear-gradient(180deg, #fffdf9, #faf4ed);
      border: 1px solid #eee4d8;
      padding: 8px 8px 2px;
    }}
    .chart-viewport {{
      width: 100%;
      height: 310px;
    }}
    .section-header {{
      display: flex;
      justify-content: space-between;
      align-items: end;
      gap: 18px;
      margin: 18px 0 18px;
    }}
    .section-header h2 {{
      margin: 0;
      color: var(--md-primary-dark);
      font-size: 1.35rem;
    }}
    .section-header p {{
      margin: 6px 0 0;
      color: var(--md-muted);
      max-width: 760px;
      line-height: 1.6;
    }}
    ul {{
      margin: 0;
      padding-left: 20px;
      line-height: 1.6;
    }}
    .modal {{
      position: fixed;
      inset: 0;
      display: none;
      align-items: center;
      justify-content: center;
      background: rgba(22, 24, 31, 0.58);
      padding: 24px;
      z-index: 30;
    }}
    .modal.open {{
      display: flex;
    }}
    .modal-card {{
      width: min(1120px, 96vw);
      max-height: 92vh;
      overflow: auto;
      background: #fff;
      border-radius: 28px;
      box-shadow: 0 28px 50px rgba(22, 24, 31, 0.24);
      padding: 22px;
    }}
    .modal-top {{
      display: flex;
      justify-content: space-between;
      align-items: center;
      gap: 16px;
      margin-bottom: 12px;
    }}
    .modal-top h3 {{
      margin: 0;
      color: var(--md-primary-dark);
      font-size: 1.25rem;
    }}
    .modal-top p {{
      margin: 6px 0 0;
      color: var(--md-muted);
    }}
    .close-button {{
      border: none;
      background: #efe5d8;
      color: var(--md-primary-dark);
      border-radius: 999px;
      width: 40px;
      height: 40px;
      font-size: 1.1rem;
      cursor: pointer;
    }}
    #modalChart {{
      width: 100%;
      height: 70vh;
      min-height: 420px;
    }}
    @media (max-width: 880px) {{
      .hero {{
        grid-template-columns: 1fr;
      }}
      .two-col {{
        grid-template-columns: 1fr;
      }}
      .hero-copy h1 {{
        font-size: 2rem;
      }}
      .chart-card-top {{
        flex-direction: column;
      }}
      .modal-card {{
        padding: 16px;
      }}
    }}
  </style>
</head>
<body>
  <div class="app-bar">
    <div class="brand">Retail Analytics Studio</div>
    <div class="chip-row">
      <span class="chip">Retail performance</span>
      <span class="chip">Sales and customer insights</span>
      <span class="chip">Forecast and retention view</span>
    </div>
  </div>
  <header>
    <section class="hero">
      <div class="hero-panel hero-copy">
        <h1>E-Commerce Business Analytics Dashboard</h1>
        <p>
          One year of UK wholesale e-commerce transactions analyzed across sales,
          products, customers, geography, timing, baskets, pricing, returns,
          cohorts, and forecasting.
        </p>
      </div>
      <div class="hero-panel hero-side">
        <div>
          <h2>Business focus</h2>
          <p>Explore the main revenue drivers, strongest customer segments, timing patterns, return risks, and the markets with the clearest expansion potential.</p>
        </div>
        <div class="chip-row">
          <span class="chip">Peak Month: {payload["summary_bits"]["top_month"]}</span>
          <span class="chip">UK Share: {format_pct(bits["uk_share"])}</span>
          <span class="chip">Best Hour: {bits["best_hour"]:02d}:00</span>
        </div>
      </div>
    </section>
  </header>
  <main>
    <section class="cards">{cards}</section>

    <section class="two-col">
      <div class="panel">
        <h2>Key Takeaways</h2>
        <ul>{finding_items}</ul>
      </div>
      <div class="panel">
        <h2>What We Recommend</h2>
        <ul>{recommendation_items}</ul>
      </div>
    </section>

    <section class="two-col">
      <div class="panel">
        <h2>Top Revenue Products</h2>
        {top_products_html}
      </div>
      <div class="panel">
        <h2>Priority Markets</h2>
        <p>The UK remains the core engine at {format_pct(bits["uk_share"])}, but international revenue already contributes {format_pct(bits["international_share"])}.</p>
        {countries_html}
      </div>
    </section>

    <section class="section-header">
      <div>
        <h2>Analysis Overview</h2>
        <p>
          The analysis below brings together sales momentum, product mix, customer retention,
          market concentration, order behavior, pricing dynamics, returns, and baseline forecasting.
        </p>
      </div>
    </section>
    <section class="chart-grid">
      {chart_blocks}
    </section>
  </main>
  <div class="modal" id="chartModal" aria-hidden="true">
    <div class="modal-card">
      <div class="modal-top">
        <div>
          <h3 id="modalTitle">Chart Detail</h3>
          <p id="modalInsight"></p>
        </div>
        <button class="close-button" id="closeModal" type="button" aria-label="Close chart">×</button>
      </div>
      <div id="modalChart"></div>
    </div>
  </div>
  <script>{plotly_js}</script>
  <script>
    const chartSpecs = {chart_json};
    const chartMeta = {{
      {", ".join(f'"{spec["id"]}": {json.dumps({"title": spec["title"], "insight": spec["insight"]})}' for spec in chart_specs)}
    }};
    const baseConfig = {{
      responsive: true,
      displaylogo: false,
      modeBarButtonsToRemove: ["lasso2d", "select2d"],
      toImageButtonOptions: {{ format: "png", filename: "ecommerce-chart", scale: 2 }}
    }};

    function renderChart(targetId, figure, height) {{
      const data = figure.data;
      const layout = Object.assign({{}}, figure.layout, {{ height }});
      Plotly.newPlot(targetId, data, layout, baseConfig);
    }}

    Object.entries(chartSpecs).forEach(([id, figure]) => {{
      renderChart(id, figure, 310);
    }});

    const modal = document.getElementById("chartModal");
    const modalTitle = document.getElementById("modalTitle");
    const modalInsight = document.getElementById("modalInsight");
    const closeModal = document.getElementById("closeModal");

    function openChart(id) {{
      const meta = chartMeta[id];
      modalTitle.textContent = meta.title;
      modalInsight.textContent = meta.insight;
      modal.classList.add("open");
      modal.setAttribute("aria-hidden", "false");
      renderChart("modalChart", chartSpecs[id], Math.min(window.innerHeight * 0.72, 720));
      setTimeout(() => Plotly.Plots.resize("modalChart"), 60);
    }}

    function hideModal() {{
      modal.classList.remove("open");
      modal.setAttribute("aria-hidden", "true");
      Plotly.purge("modalChart");
    }}

    document.querySelectorAll(".interactive-card").forEach((card) => {{
      card.addEventListener("click", (event) => {{
        if (event.target.closest(".expand-button")) return;
        openChart(card.dataset.chartId);
      }});
    }});

    document.querySelectorAll(".expand-button").forEach((button) => {{
      button.addEventListener("click", (event) => {{
        event.stopPropagation();
        openChart(button.dataset.chartId);
      }});
    }});

    closeModal.addEventListener("click", hideModal);
    modal.addEventListener("click", (event) => {{
      if (event.target === modal) hideModal();
    }});
    window.addEventListener("keydown", (event) => {{
      if (event.key === "Escape" && modal.classList.contains("open")) hideModal();
    }});
  </script>
</body>
</html>
"""
    path = DASHBOARD / "ecommerce_dashboard.html"
    path.write_text(html, encoding="utf-8")
    return path


def build_executive_summary(payload: dict) -> Path:
    h = payload["headline"]
    bits = payload["summary_bits"]
    monthly = payload["monthly"]
    path = DELIVERABLES / "executive_summary.md"
    text = f"""# Executive Summary

## Business Snapshot

- Revenue: {format_currency(h["total_revenue"])}
- Orders: {h["total_orders"]:,}
- Customers: {h["total_customers"]:,}
- Average order value: {format_currency(h["avg_order_value"])}
- Countries served: {h["countries"]}
- Returns impact: {format_currency(abs(payload["returns"]["financial_impact"]))} ({format_pct(bits["return_value_ratio"])})

## Key Findings

"""
    text += "\n".join(f"- {item}" for item in payload["findings"])
    text += f"""

## Business Recommendations

"""
    text += "\n".join(f"- {item}" for item in payload["recommendations"])
    text += f"""

## Forecast

- Baseline forecast for 2011-12: {format_currency(payload["forecast"].iloc[-3]['Revenue'])}
- Baseline forecast for 2012-01: {format_currency(payload["forecast"].iloc[-2]['Revenue'])}
- Baseline forecast for 2012-02: {format_currency(payload["forecast"].iloc[-1]['Revenue'])}

## Notes

- Final observed month ({monthly.iloc[-1]['YearMonth']}) is partial because the source data ends on December 9, 2011.
- Saturday has no recorded transactions in the source dataset.
"""
    path.write_text(text, encoding="utf-8")
    return path


def build_report_docx(payload: dict, figures: dict[str, Path]) -> Path:
    doc = Document()
    set_doc_defaults(doc)
    add_doc_title(
        doc,
        "E-Commerce Business Analytics Technical Report",
        "Online Retail dataset, Dec 2010 to Dec 2011 | Methodology, findings, and recommendations",
    )

    doc.add_heading("1. Objective", level=1)
    doc.add_paragraph(
        "This project analyzes one year of e-commerce transactions from a UK-based wholesale retailer to identify revenue drivers, customer behavior patterns, timing effects, return risks, and near-term growth opportunities."
    )

    doc.add_heading("2. Data and Preparation", level=1)
    doc.add_paragraph(
        "The source workbook contains 541,909 raw transaction rows. The project pipeline removed exact duplicates, preserved guest orders with flags, split returns into a dedicated dataset, filtered non-product adjustments from clean sales, added time features, and generated customer and product summary tables."
    )
    add_bullets(
        doc,
        [
            "5,268 exact duplicate lines removed.",
            "536,641 cleaned master rows retained for auditability.",
            "522,537 canonical sales rows used for revenue-focused analyses.",
            "10,587 return or cancellation rows isolated for return analysis.",
            "Guests retained for sales analysis but excluded from customer-level segmentation and cohorts.",
        ],
    )

    doc.add_heading("3. Headline Performance", level=1)
    doc.add_paragraph(
        f"Clean sales generated {format_currency(payload['headline']['total_revenue'])} across {payload['headline']['total_orders']:,} orders and {payload['headline']['total_customers']:,} identified customers. Average order value was {format_currency(payload['headline']['avg_order_value'])}, with a median of {format_currency(payload['headline']['median_order_value'])}."
    )
    add_picture(doc, figures["monthly_revenue"])
    doc.add_paragraph(
        "Revenue accelerated strongly in the final quarter, peaking in November 2011. September through November formed the clearest expansion period, suggesting holiday demand and merchandising momentum."
    )

    doc.add_heading("4. Product and Pricing Performance", level=1)
    doc.add_paragraph(
        "Revenue concentration is led by a small set of gift-oriented products. The leading SKU, REGENCY CAKESTAND 3 TIER, delivered the highest revenue contribution, while several low-price, high-volume items acted as basket builders."
    )
    add_picture(doc, figures["top_products_revenue"])
    add_bullets(
        doc,
        [
            "Top revenue SKUs combine broad customer reach with premium-enough pricing to preserve margin potential.",
            "The GBP 2-5 price band drives the largest share of revenue, while GBP 1-2 provides major unit volume support.",
            "Extremely high-volume single-order anomalies should be separated from ordinary merchandising decisions when planning inventory.",
        ],
    )

    doc.add_heading("5. Customer, Geography, and Retention", level=1)
    doc.add_paragraph(
        f"Repeat customers represent {format_pct(payload['summary_bits']['repeat_share'])} of identified customers. RFM segmentation shows a large Champions segment, but also a substantial At Risk / Lapsed group that should be actively reactivated."
    )
    add_picture(doc, figures["rfm_segments"])
    add_picture(doc, figures["cohort_heatmap"])
    doc.add_paragraph(
        "Geographically, the UK contributes the majority of revenue, but international markets such as the Netherlands, EIRE, Germany, and France offer strong expansion efficiency because they already generate meaningful value from relatively few customers."
    )

    doc.add_heading("6. Timing, Returns, and Forecast", level=1)
    doc.add_paragraph(
        f"Value peaks during business hours, with the strongest hours centered from 10:00 to 15:00 and particularly strong trading on {payload['summary_bits']['best_day']}. Returns and cancellations reduced value by {format_currency(abs(payload['returns']['financial_impact']))}, equivalent to {format_pct(payload['summary_bits']['return_value_ratio'])} of clean sales."
    )
    add_picture(doc, figures["hour_of_day"])
    add_picture(doc, figures["forecast"])
    doc.add_paragraph(
        "A simple linear-trend baseline suggests continued revenue strength into early 2012, although the model should be treated as directional because it does not yet incorporate holidays, promotions, or inventory constraints."
    )

    doc.add_heading("7. Recommendations", level=1)
    add_bullets(doc, payload["recommendations"])

    doc.add_heading("8. Limitations", level=1)
    add_bullets(
        doc,
        [
            "The source data covers only one year, which limits long-run trend estimation.",
            "The final month is partial because the dataset ends on December 9, 2011.",
            "Some return rows reflect manual adjustments or operational corrections rather than merchandise dissatisfaction.",
            "Profitability cannot be estimated directly because product cost data is not available.",
        ],
    )

    path = REPORTS / "ecommerce_technical_report.docx"
    doc.save(path)
    return path


def add_slide_title(slide, title: str, subtitle: str | None = None) -> None:
    title_box = slide.shapes.add_textbox(PptInches(0.6), PptInches(0.3), PptInches(11.8), PptInches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = PptPt(24)
    p.font.bold = True
    p.font.color.rgb = PptRGBColor(20, 58, 82)
    if subtitle:
        sub_box = slide.shapes.add_textbox(PptInches(0.6), PptInches(0.95), PptInches(11.2), PptInches(0.4))
        p2 = sub_box.text_frame.paragraphs[0]
        p2.text = subtitle
        p2.font.size = PptPt(11)
        p2.font.color.rgb = PptRGBColor(84, 101, 117)


def add_bullet_list(slide, items: list[str], left: float, top: float, width: float, height: float, font_size: int = 18) -> None:
    box = slide.shapes.add_textbox(PptInches(left), PptInches(top), PptInches(width), PptInches(height))
    tf = box.text_frame
    tf.word_wrap = True
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = PptPt(font_size)
        p.font.color.rgb = PptRGBColor(31, 41, 51)
        p.bullet = True


def add_picture_slide(slide, image_path: Path, left: float, top: float, width: float) -> None:
    slide.shapes.add_picture(str(image_path), PptInches(left), PptInches(top), width=PptInches(width))


def add_metric_banner(slide, metrics: list[tuple[str, str]]) -> None:
    start_left = 0.6
    width = 2.85
    gap = 0.18
    for idx, (label, value) in enumerate(metrics):
        shape = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            PptInches(start_left + idx * (width + gap)),
            PptInches(1.35),
            PptInches(width),
            PptInches(1.0),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = PptRGBColor(238, 243, 247)
        shape.line.color.rgb = PptRGBColor(220, 228, 235)
        tf = shape.text_frame
        tf.clear()
        p1 = tf.paragraphs[0]
        p1.text = label
        p1.font.size = PptPt(10)
        p1.font.color.rgb = PptRGBColor(92, 111, 127)
        p1.alignment = PP_ALIGN.CENTER
        p2 = tf.add_paragraph()
        p2.text = value
        p2.font.size = PptPt(18)
        p2.font.bold = True
        p2.font.color.rgb = PptRGBColor(20, 58, 82)
        p2.alignment = PP_ALIGN.CENTER


def build_pptx(payload: dict, figures: dict[str, Path]) -> Path:
    prs = Presentation()
    prs.slide_width = PptInches(13.333)
    prs.slide_height = PptInches(7.5)
    blank = prs.slide_layouts[6]

    # Slide 1
    slide = prs.slides.add_slide(blank)
    bg = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = PptRGBColor(20, 58, 82)
    bg.line.fill.background()
    title = slide.shapes.add_textbox(PptInches(0.8), PptInches(0.9), PptInches(8.4), PptInches(1.2))
    p = title.text_frame.paragraphs[0]
    p.text = "E-Commerce Business Analytics"
    p.font.size = PptPt(28)
    p.font.bold = True
    p.font.color.rgb = PptRGBColor(255, 255, 255)
    sub = slide.shapes.add_textbox(PptInches(0.8), PptInches(2.0), PptInches(6.5), PptInches(1.0))
    p2 = sub.text_frame.paragraphs[0]
    p2.text = "Executive findings and recommendations\nOnline Retail | Dec 2010 - Dec 2011"
    p2.font.size = PptPt(16)
    p2.font.color.rgb = PptRGBColor(222, 235, 245)
    add_metric_banner(
        slide,
        [
            ("Revenue", format_currency(payload["headline"]["total_revenue"])),
            ("Orders", f"{payload['headline']['total_orders']:,}"),
            ("Customers", f"{payload['headline']['total_customers']:,}"),
            ("Countries", str(payload["headline"]["countries"])),
        ],
    )

    # Slide 2
    slide = prs.slides.add_slide(blank)
    add_slide_title(slide, "What mattered most")
    add_bullet_list(slide, payload["findings"], 0.8, 1.5, 11.2, 4.8)

    # Slide 3
    slide = prs.slides.add_slide(blank)
    add_slide_title(slide, "Sales overview", "Revenue expanded sharply in the final quarter.")
    add_picture_slide(slide, figures["monthly_revenue"], 0.7, 1.3, 7.1)
    add_bullet_list(
        slide,
        [
            f"Peak month: {payload['summary_bits']['top_month']} at {format_currency(payload['summary_bits']['top_month_revenue'])}.",
            "September to November created the strongest revenue run-rate.",
            "Average order value reached GBP 518, confirming wholesale-sized baskets.",
        ],
        8.1,
        1.55,
        4.4,
        3.0,
        16,
    )

    # Slide 4
    slide = prs.slides.add_slide(blank)
    add_slide_title(slide, "Products and pricing", "A small set of hero SKUs carries outsized value.")
    add_picture_slide(slide, figures["top_products_revenue"], 0.7, 1.4, 6.2)
    add_picture_slide(slide, figures["price_bands"], 7.2, 1.4, 5.3)

    # Slide 5
    slide = prs.slides.add_slide(blank)
    add_slide_title(slide, "Customer behavior", "Repeat business is meaningful, but lapse risk is large.")
    add_picture_slide(slide, figures["rfm_segments"], 0.7, 1.4, 6.0)
    add_bullet_list(
        slide,
        [
            f"Repeat customers represent {format_pct(payload['summary_bits']['repeat_share'])} of identified customers.",
            "Champions and Loyal customers justify VIP retention motions.",
            "At Risk / Lapsed is the largest segment and the biggest recovery opportunity.",
        ],
        7.1,
        1.55,
        5.1,
        3.3,
        16,
    )

    # Slide 6
    slide = prs.slides.add_slide(blank)
    add_slide_title(slide, "Geographic concentration", "The UK dominates, but several export markets already prove demand.")
    add_picture_slide(slide, figures["top_countries"], 0.7, 1.35, 7.0)
    add_bullet_list(
        slide,
        [
            f"UK share of revenue: {format_pct(payload['summary_bits']['uk_share'])}.",
            "Netherlands, EIRE, Germany, and France are the best initial expansion markets.",
            "International growth can be selective rather than broad and expensive.",
        ],
        8.0,
        1.55,
        4.4,
        3.2,
        16,
    )

    # Slide 7
    slide = prs.slides.add_slide(blank)
    add_slide_title(slide, "When to promote", "High-value trading clusters around midweek business hours.")
    add_picture_slide(slide, figures["day_of_week"], 0.7, 1.45, 5.7)
    add_picture_slide(slide, figures["hour_of_day"], 6.7, 1.45, 5.7)

    # Slide 8
    slide = prs.slides.add_slide(blank)
    add_slide_title(slide, "Basket economics", "Order values are healthy but skewed by a long right tail.")
    add_picture_slide(slide, figures["basket_distribution"], 0.8, 1.5, 6.7)
    add_bullet_list(
        slide,
        [
            f"Average basket value: {format_currency(payload['basket_metrics']['avg_value'])}.",
            f"Median basket value: {format_currency(payload['basket_metrics']['median_value'])}.",
            f"Average distinct items per order: {payload['basket_metrics']['avg_items']:.1f}.",
        ],
        8.0,
        1.7,
        4.2,
        2.5,
        16,
    )

    # Slide 9
    slide = prs.slides.add_slide(blank)
    add_slide_title(slide, "Returns and operational leakage", "Returns are manageable overall but concentrated in a handful of codes.")
    add_picture_slide(slide, figures["top_returns"], 0.7, 1.5, 6.2)
    add_bullet_list(
        slide,
        [
            f"Returns impact: {format_currency(abs(payload['returns']['financial_impact']))}.",
            f"Return-line rate vs total lines: {format_pct(payload['returns']['return_rate_vs_sales'])}.",
            "Manual and operational adjustment codes should be separated from real merchandise return root causes.",
        ],
        7.4,
        1.7,
        4.8,
        2.8,
        16,
    )

    # Slide 10
    slide = prs.slides.add_slide(blank)
    add_slide_title(slide, "Retention and customer lifetime", "Early repeat behavior determines long-run customer value.")
    add_picture_slide(slide, figures["cohort_heatmap"], 0.7, 1.45, 11.8)

    # Slide 11
    slide = prs.slides.add_slide(blank)
    add_slide_title(slide, "Forecast", "Baseline trend points to continued growth into early 2012.")
    add_picture_slide(slide, figures["forecast"], 0.8, 1.55, 6.6)
    add_bullet_list(
        slide,
        [
            f"Forecast 2012-01: {format_currency(payload['forecast'].iloc[-2]['Revenue'])}.",
            f"Forecast 2012-02: {format_currency(payload['forecast'].iloc[-1]['Revenue'])}.",
            "Treat as a baseline trend, then refine with promotion and seasonality assumptions.",
        ],
        8.0,
        1.75,
        4.2,
        2.8,
        16,
    )

    # Slide 12
    slide = prs.slides.add_slide(blank)
    add_slide_title(slide, "Recommended actions", "What the business should do next")
    add_bullet_list(slide, payload["recommendations"], 0.8, 1.5, 11.2, 4.8)

    path = PRESENTATION / "ecommerce_business_analytics_deck.pptx"
    prs.save(path)
    return path


def build_metadata(payload: dict, paths: dict[str, Path]) -> Path:
    meta = {
        "headline_metrics": payload["headline"],
        "findings": payload["findings"],
        "recommendations": payload["recommendations"],
        "artifacts": {key: str(path.relative_to(ROOT)) for key, path in paths.items()},
    }
    path = DELIVERABLES / "deliverables_manifest.json"
    path.write_text(json.dumps(meta, indent=2, default=str), encoding="utf-8")
    return path


def main() -> None:
    ensure_dirs()
    setup_style()
    payload = build_analysis_payload()
    figures = build_figures(payload)
    build_tables(payload)
    dashboard_path = build_dashboard(payload, figures)
    summary_path = build_executive_summary(payload)
    report_path = build_report_docx(payload, figures)
    deck_path = build_pptx(payload, figures)
    manifest_path = build_metadata(
        payload,
        {
            "dashboard": dashboard_path,
            "executive_summary": summary_path,
            "technical_report": report_path,
            "presentation": deck_path,
        },
    )

    print("Deliverables generated:")
    for path in [
        dashboard_path,
        summary_path,
        report_path,
        deck_path,
        manifest_path,
    ]:
        print(f"  - {path.relative_to(ROOT)}")


if __name__ == "__main__":
    main()
